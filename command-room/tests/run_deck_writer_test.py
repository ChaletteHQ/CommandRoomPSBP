#!/usr/bin/env python3
"""
Tests for `shared/scripts/deck_writer.py` (SPEC OUT6 — the board-pack .pptx
chokepoint) + the visual_gate pptx→PDF rung dispatch.

Covers the contract deck_writer must satisfy:

  - GRAMMAR pins match the prose table in shared/DECK_GRAMMAR.md (machine copy
    and prose copy can never drift silently)
  - section→slide map: exec section → title slide only (no agenda slide);
    kpi / wins / concerns / decisions / asks / appendix treatments
  - one message per slide (a slide never carries two section headings)
  - caps: wins/concerns 5 rows + honest overflow note; asks 3 (refusal, not a
    squeeze); appendix ≤6 bullets/slide, ≤2 slides/section, continuation slide
  - drop-empty (F-60): "(nothing logged)" sections and all-placeholder
    sections contribute no slide; an all-empty deck REFUSES (title-only deck
    never ships)
  - placeholder mixed into real content refuses (DeckGrammarError)
  - empty tile refused at the chokepoint (components.validate_tiles)
  - leak scan BEFORE save, proven with a poisoned fixture — DeckLeakError and
    NO file on disk
  - brand resolution incl. per-org override (accent hex lands in slide XML);
    unconfigured = byte-stable defaults
  - deterministic: same sections + same brand → identical .pptx zip payload
    (per-entry bytes; raw archive bytes carry 2s-resolution zip timestamps)
  - docx/deck parity: the deck's KPI values / wins / asks match the .docx
    verbatim (same assembled sections, two renderers)
  - the install path: pip failure → DeckDependencyError (honest stop), never
    a hand-built fallback
  - visual_gate: .pptx paths take _PPTX_TO_PDF_LADDER, .docx paths still take
    _DOCX_TO_PDF_LADDER (the pinned docx ladder is undisturbed)

python-pptx self-installs on first make_deck use (PYTHON_PPTX_PIN). If it is
genuinely absent AND uninstallable here, render-path checks SKIP LOUDLY (the
plan/grammar checks still run — they are pptx-free by design); the install
honest-stop check runs regardless (it mocks the dependency away).

No hardcoded future dates (G14): fixtures carry no dates the writer computes
with; the one date literal is rendered verbatim as table text.
"""
import io
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "shared", "scripts"))

import deck_writer
from deck_writer import (
    GRAMMAR,
    DeckDependencyError,
    DeckGrammarError,
    DeckLeakError,
    build_slide_plan,
    make_deck,
)

results = {"pass": 0, "fail": 0, "failures": []}

# Consolidated 2026-07-19 (FB bundle): the local copy that mirrored
# run_charts_test.py while out7-kpi-scorecard was in flight is now the shared
# tests/ooxml_payload_lib.py helper (OUT7 merged @ c122137). Both suites import it.
from ooxml_payload_lib import zip_payload_identical as _zip_payload_identical  # noqa: E402


def check(name, condition, detail=""):
    if condition:
        results["pass"] += 1
        print(f"  PASS  {name}")
    else:
        results["fail"] += 1
        results["failures"].append(f"{name} ({detail})" if detail else name)
        print(f"  FAIL  {name}{(' — ' + detail) if detail else ''}")


# Real-substrate section shapes (per the real-data fixture gotcha): tile band
# with target-delta arrows, KPI table with headers, quantify-tagged wins,
# decisions table, capped asks, an overflowing appendix, a "(nothing logged)"
# section, and the docx's sanctioned all-placeholder asks form.
SECTIONS = [
    {"heading": "Executive summary",
     "bullets": ["MRR $478K, +13% MoM", "Closed Sample Co — $52K"]},
    {"heading": "KPIs vs targets",
     "tiles": [{"label": "MRR ▲", "value": "$478K"},
               {"label": "NRR ▲", "value": "134%"},
               {"label": "Runway", "value": "18.4 mo"}],
     "table": {"headers": ["Metric", "Current", "Target", "vs Target"],
               "rows": [["MRR", "$478K", "$470K", "+1.7%"],
                        ["NRR", "134%", "125%", "+9 pts"]]}},
    {"heading": "Wins",
     "bullets": ["Closed Sample Co — $52K", "Hit ARR milestone in 60 days",
                 "Pilot renewed", "Fourth win landed", "Fifth win landed",
                 "Sixth win overflows", "Seventh win overflows"]},
    {"heading": "Concerns",
     "bullets": ["Lost Northstar Sample (timing); postmortem attached"]},
    {"heading": "Decisions logged this period",
     "table": {"headers": ["Decision", "Owner", "Date"],
               "rows": [["Paused consumer pilot", "CEO", "2026-05-04"]]}},
    {"heading": "Asks for the board",
     "bullets": ["Extend Q3 hiring envelope — $340K",
                 "Intro to the operator network"]},
    {"heading": "Appendix — pipeline by stage",
     "bullets": [f"Stage line {i}" for i in range(1, 15)]},  # 14 → 6+6, +2 note
    {"heading": "Appendix — hiring slate",
     "body": "(nothing logged this period)"},
]
KW = dict(title="fallback title", subtitle="Sample Org — reporting period",
          exec_header={"verdict": "MRR up 13% against a 10% target"})


def _plan():
    return build_slide_plan(SECTIONS, **KW)


def test_grammar_pins_match_prose():
    md = Path(HERE, "..", "shared", "DECK_GRAMMAR.md").read_text(encoding="utf-8")
    for pin, value in GRAMMAR.items():
        if pin in ("slide_w_in", "slide_h_in"):
            continue  # asserted together below
        found = re.search(rf"`{pin}`\s*\|\s*(\d+)", md)
        check(f"grammar pin {pin} in prose", found is not None
              and int(found.group(1)) == value,
              f"prose={found and found.group(1)} code={value}")
    check("grammar pin slide dims in prose",
          f"{GRAMMAR['slide_w_in']} / {GRAMMAR['slide_h_in']}" in md)
    floor = GRAMMAR["font_floor_pt"]
    check("font floor is structural (sizes table)",
          min(deck_writer._SIZES_PT.values()) >= floor)


def test_section_slide_map():
    plan = _plan()
    kinds = [s["slide"] for s in plan]
    check("title slide first", kinds[0] == "title")
    check("verdict is the headline", plan[0]["title"] == KW["exec_header"]["verdict"])
    check("no agenda slide / exec section dropped",
          not any("Executive" in str(s.get("heading")) for s in plan))
    check("kpi slide present with tiles+table",
          any(s["slide"] == "kpi" and s.get("tiles") and s.get("table") for s in plan))
    wins = next(s for s in plan if s.get("heading") == "Wins")
    check("wins capped at 5 rows", len(wins["rows"]) == GRAMMAR["max_rows_wins_concerns"])
    check("wins overflow is an honest note", wins["note"] == "+2 more in the full pack")
    check("quantify tag split", wins["rows"][0]["tag"] == "$52K"
          and wins["rows"][0]["text"] == "Closed Sample Co")
    check("decisions render as a table",
          any(s["slide"] == "table" and "Decisions" in s["heading"] for s in plan))
    asks = next(s for s in plan if "Asks" in str(s.get("heading")))
    check("asks slide carries both asks with tags",
          len(asks["rows"]) == 2 and asks["rows"][0]["tag"] == "$340K")
    appendix = [s for s in plan if "pipeline" in str(s.get("heading"))]
    check("appendix ≤2 slides with continuation",
          len(appendix) == GRAMMAR["max_content_slides_per_section"]
          and appendix[1]["heading"].endswith("(cont.)"))
    check("appendix bullets per slide ≤ cap",
          all(len(s["rows"]) <= GRAMMAR["max_bullets_per_slide"] for s in appendix))
    check("appendix overflow note past slide cap",
          appendix[-1]["note"] == "+2 more in the full pack")
    check("(nothing logged) section dropped",
          not any("hiring" in str(s.get("heading", "")) for s in plan))
    check("one message per slide",
          all(isinstance(s.get("heading"), str) or s["slide"] == "title" for s in plan))
    # Second-eyes F-1: the exec classifier is word-bounded. "Execution risks"
    # is a CONTENT section (an unbounded \bexec prefix-match silently dropped
    # it from the deck while the docx carried it); "Exec summary" still
    # classifies exec and gets no content slide.
    plan_exec = build_slide_plan(
        [{"heading": "Wins", "bullets": ["Real win logged"]},
         {"heading": "Execution risks", "bullets": ["Vendor slip risk"]},
         {"heading": "Exec summary", "bullets": ["docx page-1 bullet"]}], **KW)
    check("'Execution risks' renders as a content slide (F-1)",
          any(s.get("heading") == "Execution risks" for s in plan_exec))
    check("'Exec summary' still title-slide-only (F-1)",
          not any("Exec summary" in str(s.get("heading", "")) for s in plan_exec))


def test_refusals():
    try:
        build_slide_plan(SECTIONS, kind="memo", **KW)
        check("scope fence: non-board_pack kind refused", False)
    except ValueError as e:
        check("scope fence: non-board_pack kind refused", "scope fence" in str(e))
    try:
        build_slide_plan([{"heading": "Wins", "bullets": ["real", "[insert win here]"]}], **KW)
        check("placeholder mixed into real content refused", False)
    except DeckGrammarError:
        check("placeholder mixed into real content refused", True)
    try:
        build_slide_plan(
            [{"heading": "Asks", "bullets": ["a1 — $1K", "a2", "a3", "a4"]}], **KW)
        check("four asks refused (EXEC1 cap)", False)
    except DeckGrammarError as e:
        check("four asks refused (EXEC1 cap)", "cap" in str(e))
    try:
        build_slide_plan(SECTIONS, title="t", subtitle="s",
                         asks=[{"text": f"a{i}"} for i in range(4)])
        check("four asks via kwarg refused", False)
    except DeckGrammarError:
        check("four asks via kwarg refused", True)
    try:
        build_slide_plan([{"heading": "Wins", "body": "(nothing logged this period)"}], **KW)
        check("all-empty deck refused (no title-only deck)", False)
    except DeckGrammarError as e:
        check("all-empty deck refused (no title-only deck)", "empty" in str(e))
    try:
        build_slide_plan(
            [{"heading": "KPIs", "tiles": [{"label": "MRR", "value": ""}],
              "table": {"rows": [["x"]]}}], **KW)
        check("empty tile refused at chokepoint", False)
    except ValueError as e:
        check("empty tile refused at chokepoint", "DROPPED" in str(e))
    # The docx's sanctioned all-placeholder asks form drops (never renders).
    plan = build_slide_plan(
        SECTIONS[:2] + [{"heading": "Asks for the board", "body": "[add asks here]"}], **KW)
    check("all-placeholder asks section drops, deck still renders",
          not any("Asks" in str(s.get("heading", "")) for s in plan) and len(plan) >= 2)


def _pptx_available():
    try:
        import pptx  # noqa: F401
        return True
    except ImportError:
        return False


def test_render_and_brand_and_determinism(tmp):
    if not _pptx_available():
        # Loud skip — the self-install normally makes this unreachable.
        print("  SKIP  render checks — python-pptx absent and not installed "
              "(self-install must have failed; the honest-stop check below still runs)")
        return
    from pptx import Presentation

    out1 = os.path.join(tmp, "deck1.pptx")
    out2 = os.path.join(tmp, "deck2.pptx")
    make_deck(out1, SECTIONS, **KW)
    make_deck(out2, SECTIONS, **KW)
    check("deck saved non-empty", os.path.getsize(out1) > 10_000)
    check("deterministic payload for fixed input",
          _zip_payload_identical(out1, out2),
          "payload-identity couples to python-pptx internals (part naming, "
          "part ordering, XML serialization) and is only promised WITHIN the "
          f"pinned version (PYTHON_PPTX_PIN={deck_writer.PYTHON_PPTX_PIN}). "
          "If this failed right after a pin bump, payload drift is EXPECTED — "
          "verify the slide XML is still deterministic, then re-baseline "
          "deliberately; do not chase phantom nondeterminism.")

    prs = Presentation(out1)
    n_plan = len(_plan())
    check("slide count matches plan", len(prs.slides._sldIdLst) == n_plan,
          f"slides={len(prs.slides._sldIdLst)} plan={n_plan}")
    xml = "".join(s._element.xml for s in prs.slides)
    from brand import DEFAULT_BRAND
    check("default brand accent in slide XML",
          DEFAULT_BRAND["palette"]["accent"] in xml)
    check("default heading font in slide XML",
          DEFAULT_BRAND["fonts"]["heading"] in xml)

    # Per-org override: org accent must land in the deck (board decks are
    # per-org documents). Placeholder org only (org-name leak gotcha).
    entities = {"workspace": {}, "entities": {"orgs": [
        {"id": "org_sample", "canonical_name": "Sample Org",
         "brand": {"palette": {"accent": "8A5A2B"}}}]}}
    from brand import get_brand
    org_brand = get_brand(entities, org_id="org_sample")
    out3 = os.path.join(tmp, "deck3.pptx")
    make_deck(out3, SECTIONS, brand=org_brand, **KW)
    xml3 = "".join(s._element.xml for s in Presentation(out3).slides)
    check("org brand override accent in slide XML", "8A5A2B" in xml3)
    # not read_bytes() != : zip-timestamp drift alone would make raw bytes
    # differ and false-pass this even if the override never landed.
    check("org override changes payload vs default",
          not _zip_payload_identical(out3, out1))


def test_leak_scan_before_save(tmp):
    poisoned = [{"heading": "Wins",
                 "bullets": ["project_020 shipped — $52K",
                             "Numbers live in events.jsonl now"]}]
    out = os.path.join(tmp, "poisoned.pptx")
    try:
        make_deck(out, poisoned, **KW)
        check("poisoned fixture refused", False)
    except DeckLeakError as e:
        check("poisoned fixture refused", "project_020" in str(e))
    check("no file written on leak refusal", not os.path.exists(out))

    # Second-eyes F-2: the brand footer line is painted on every slide but
    # does not pass through the plan — it must be scanned too. A poisoned
    # footer refuses the render with no file on disk.
    from brand import DEFAULT_BRAND
    poisoned_brand = {**DEFAULT_BRAND,
                      "footer_line": "Command Room — see events.jsonl"}
    out2 = os.path.join(tmp, "poisoned_footer.pptx")
    try:
        make_deck(out2, SECTIONS, brand=poisoned_brand, **KW)
        check("poisoned brand footer_line refused (F-2)", False)
    except DeckLeakError as e:
        check("poisoned brand footer_line refused (F-2)",
              "events.jsonl" in str(e))
    check("no file written on footer leak refusal (F-2)",
          not os.path.exists(out2))


def test_parity_with_docx(tmp):
    if not _pptx_available():
        print("  SKIP  parity — python-pptx absent")
        return
    from brief_writer import make_brief
    from pptx import Presentation
    from docx import Document

    docx_path = os.path.join(tmp, "pack.docx")
    pptx_path = os.path.join(tmp, "pack.pptx")
    asks = [{"text": "Extend Q3 hiring envelope — $340K"}]
    make_brief(docx_path, brief_kind="board_pack", title=KW["title"],
               subtitle=KW["subtitle"], sections=SECTIONS,
               exec_header=KW["exec_header"], asks=asks,
               contract="off", voice_gate="off")
    make_deck(pptx_path, SECTIONS, title=KW["title"], subtitle=KW["subtitle"],
              exec_header=KW["exec_header"], asks=asks)

    docx_text = "\n".join(p.text for p in Document(docx_path).paragraphs)
    for tbl in Document(docx_path).tables:
        for row in tbl.rows:
            docx_text += "\n" + "\t".join(c.text for c in row.cells)
    pptx_text = ""
    for slide in Presentation(pptx_path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                pptx_text += "\n" + shape.text_frame.text
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    pptx_text += "\n" + "\t".join(c.text for c in row.cells)

    # Structural diff: every KPI value, win, and ask the docx carries must
    # appear in the deck verbatim (one assembly, two renderers).
    for value in ("$478K", "$470K", "134%", "125%", "18.4 mo",
                  "Closed Sample Co", "$52K",
                  "Extend Q3 hiring envelope", "$340K",
                  "Paused consumer pilot"):
        check(f"parity: {value!r} in both artifacts",
              value in docx_text and value in pptx_text,
              f"docx={value in docx_text} pptx={value in pptx_text}")


def test_install_honest_stop():
    """Mock the dependency away + break pip: make_deck must raise
    DeckDependencyError (the one-line honest stop), never fall through."""
    real_mod = sys.modules.get("pptx")
    real_run = deck_writer.subprocess.run

    def _pip_fails(*a, **k):
        raise subprocess.CalledProcessError(1, a[0] if a else "pip")

    try:
        sys.modules["pptx"] = None  # forces `import pptx` → ImportError
        deck_writer.subprocess = type(sys)("subprocess_stub")
        deck_writer.subprocess.run = _pip_fails
        try:
            make_deck(os.path.join(tempfile.gettempdir(), "never.pptx"),
                      SECTIONS, **KW)
            check("install failure → DeckDependencyError", False)
        except DeckDependencyError as e:
            msg = str(e)
            check("install failure → DeckDependencyError", True)
            check("honest-stop message is one line and names the fallback",
                  "\n" not in msg and ".docx pack is unaffected" in msg, msg)
    finally:
        deck_writer.subprocess = __import__("subprocess")
        deck_writer.subprocess.run = real_run
        if real_mod is not None:
            sys.modules["pptx"] = real_mod
        else:
            sys.modules.pop("pptx", None)


def test_visual_gate_pptx_dispatch(tmp):
    import visual_gate

    calls = {"docx": 0, "pptx": 0}

    def _rec(bucket):
        def rung(path, out_dir):
            calls[bucket] += 1
            return None
        return rung

    docx_file = os.path.join(tmp, "probe.docx")
    pptx_file = os.path.join(tmp, "probe.pptx")
    Path(docx_file).write_bytes(b"stub")
    Path(pptx_file).write_bytes(b"stub")

    env_before = os.environ.pop("CR_VISUAL_GATE", None)
    orig_docx, orig_pptx = visual_gate._DOCX_TO_PDF_LADDER, visual_gate._PPTX_TO_PDF_LADDER
    try:
        visual_gate._DOCX_TO_PDF_LADDER = (_rec("docx"),)
        visual_gate._PPTX_TO_PDF_LADDER = (_rec("pptx"),)
        out = visual_gate.render_preview(pptx_file)
        check("pptx path → pptx ladder only, None on no-renderer",
              out is None and calls == {"docx": 0, "pptx": 1}, str(calls))
        out = visual_gate.render_preview(docx_file)
        check("docx path → docx ladder undisturbed",
              out is None and calls == {"docx": 1, "pptx": 1}, str(calls))
        check("pptx ladder ends in the shared soffice rung",
              orig_pptx[-1] is visual_gate._docx_to_pdf_soffice)
    finally:
        visual_gate._DOCX_TO_PDF_LADDER = orig_docx
        visual_gate._PPTX_TO_PDF_LADDER = orig_pptx
        if env_before is not None:
            os.environ["CR_VISUAL_GATE"] = env_before


def main():
    print("deck_writer test suite (SPEC OUT6)")
    with tempfile.TemporaryDirectory() as tmp:
        test_grammar_pins_match_prose()
        test_section_slide_map()
        test_refusals()
        test_render_and_brand_and_determinism(tmp)
        test_leak_scan_before_save(tmp)
        test_parity_with_docx(tmp)
        test_install_honest_stop()
        test_visual_gate_pptx_dispatch(tmp)
    print(f"\n{results['pass']} passed, {results['fail']} failed")
    if results["fail"]:
        for f in results["failures"]:
            print(f"  FAILED: {f}")
        sys.exit(1)


if __name__ == "__main__":
    main()
