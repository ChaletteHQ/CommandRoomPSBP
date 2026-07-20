#!/usr/bin/env python3
"""
Board-pack .pptx chokepoint (SPEC OUT6).

Single source of truth for the board pack's slide companion. Mirrors the
`brief_writer.make_brief` posture for the one deck the plugin renders:
validates input shapes, resolves ALL colors/fonts through `brand.get_brand()`
(org override honored — board decks are per-org documents), runs the leak scan
on every text run BEFORE save, refuses placeholder content and drops empty
sections (F-60). The prose contract is `shared/DECK_GRAMMAR.md`; the machine
copy of every grammar number is the `GRAMMAR` dict below — change one, change
the other, same commit (`tests/run_deck_writer_test.py` asserts the pins).

ONE KIND. `make_deck` renders the board-pack deck and nothing else — other
composers getting a deck path is FUTURE_WORK behind its own spec.

DEPENDENCY POSTURE (brief_writer-style, but lazy)
-------------------------------------------------
python-pptx is pinned to PYTHON_PPTX_PIN and self-installs on FIRST USE
(inside `make_deck`), not at import — so importing this module for its grammar
constants can never die on a machine without the dependency (the FS-15 F-1
lesson: a module-level hard import in a shared helper takes unrelated readers
down with it). If the install fails, `make_deck` raises `DeckDependencyError`
with a one-line message and NO file is written — the skill says so, delivers
the .docx, and stops. A failed chokepoint never falls back to freelance deck
generation (the render_and_persist rule, generalized).

TWO-STAGE RENDER (plan → paint)
-------------------------------
`build_slide_plan(sections, ...)` is a PURE function: sections in, a list of
slide dicts out, every grammar rule applied, no python-pptx needed. That is
where the section→slide map, the caps, the drop-empty/placeholder rules, and
the overflow handling live — and it is what the leak scan runs over (every
string that could reach a text run is in the plan, plus the resolved brand
footer line, the one painted string that originates in brand config instead).
`make_deck` then paints the plan with python-pptx. Same plan + same brand →
identical slide XML.

Stdlib only at import (brand + components + docx_leak_scanner are stdlib-only
plugin modules); python-pptx on first use.
"""
from __future__ import annotations

import re
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Optional

_HERE = Path(__file__).resolve().parent
if str(_HERE) not in sys.path:
    sys.path.insert(0, str(_HERE))

from brand import get_brand  # noqa: E402
from components import validate_tiles as _validate_tiles  # noqa: E402
from components import validate_table as _validate_table  # noqa: E402
from docx_leak_scanner import scan_text_for_leaks  # noqa: E402
# SPEC OUT3 — chart validation (pure) + strings for the plan-time leak scan;
# the PNG render happens at paint time via charts.try_chart_png. Stdlib-only
# plugin module, same import category as the three above.
from charts import ChartDataError, chart_strings, validate_chart  # noqa: E402


PYTHON_PPTX_PIN = "1.0.2"

# The machine copy of shared/DECK_GRAMMAR.md's pinned numbers. The test suite
# asserts these against the prose table — change both in the same commit.
GRAMMAR = {
    "max_bullets_per_slide": 6,
    "max_rows_wins_concerns": 5,
    "max_asks": 3,                        # EXEC1 MAX_ASKS
    "max_table_rows_per_slide": 8,
    "max_content_slides_per_section": 2,
    "font_floor_pt": 12,
    "slide_w_in": 13.333,
    "slide_h_in": 7.5,
}

SUPPORTED_DECK_KINDS = frozenset({"board_pack"})


class DeckDependencyError(RuntimeError):
    """python-pptx could not be imported or installed. The caller says so in
    one line, delivers the .docx, and stops — no improvised deck."""


class DeckGrammarError(ValueError):
    """The sections payload violates the deck grammar (placeholder text mixed
    into real content, an over-cap ask list, or nothing renderable at all)."""


class DeckLeakError(RuntimeError):
    """Forbidden tokens found in the slide plan (pre-save). No file written."""


def _ensure_python_pptx() -> None:
    """Import python-pptx, self-installing the pin on first use (idempotent,
    brief_writer-style). Raises DeckDependencyError — never a bare
    ImportError — so the skill's honest-stop line has one exception to catch."""
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass
    print(
        f"Installing python-pptx (=={PYTHON_PPTX_PIN}) — one-time setup. "
        "(Plugin requires this for the board-deck .pptx companion. "
        "See README Requirements section to pre-install in locked-down environments.)",
        file=sys.stderr,
    )
    try:
        subprocess.run(
            [sys.executable, "-m", "pip", "install", "--quiet",
             f"python-pptx=={PYTHON_PPTX_PIN}"],
            check=True,
        )
        import pptx  # noqa: F401
    except Exception as e:
        raise DeckDependencyError(
            f"python-pptx (=={PYTHON_PPTX_PIN}) is not available and could not "
            f"be installed ({type(e).__name__}) — the board deck was not "
            f"rendered; the .docx pack is unaffected."
        ) from e


# ---------------------------------------------------------------------------
# Section classification (the section → slide map, DECK_GRAMMAR table)
# ---------------------------------------------------------------------------

_ROLE_PATTERNS = [
    # Word-bounded: "Executive summary" / "Exec header" classify exec (title
    # slide, no content slide); "Execution risks" must NOT — an unbounded
    # \bexec prefix-match silently dropped such sections from the deck while
    # the docx carried them (second-eyes F-1).
    ("exec", re.compile(r"\bexec(?:utive)?\b", re.I)),
    ("kpi", re.compile(r"\bkpis?\b", re.I)),
    ("wins", re.compile(r"\bwins?\b", re.I)),
    ("concerns", re.compile(r"\bconcerns?\b", re.I)),
    ("decisions", re.compile(r"\bdecisions?\b", re.I)),
    ("asks", re.compile(r"\basks?\b", re.I)),
]

# Placeholder text refuses a render when mixed into real content; a section
# that is ONLY placeholder (the docx's sanctioned "[add asks here]") drops.
_PLACEHOLDER_RE = re.compile(
    r"\[(?:add|insert|tbd|todo)[^\]]*\]|lorem ipsum|\bTODO\b", re.I
)
# The docx's validated empty-section form — contributes no slide (drop-empty).
_NOTHING_LOGGED_RE = re.compile(r"\(nothing logged", re.I)

# The EXEC1 quantify tag: a trailing money / percent / delta / duration token
# renders as an accent tag run ("Closed Sample Deal — $52K" → tag "$52K").
_QUANT_TAG_RE = re.compile(
    r"^(?P<text>.+?)\s*[—–-]\s*"
    r"(?P<tag>[+\-]?\$[\d.,]+\s?[KMBkmb]?|[+\-]?[\d.,]+\s?(?:%|pts?|bps)|"
    r"[\d.,]+\s?(?:mo|months?|weeks?|days?|FTE)s?)\s*$"
)


def _classify(heading: str) -> str:
    for role, pat in _ROLE_PATTERNS:
        if pat.search(str(heading or "")):
            return role
    return "content"


def _split_quantify_tag(text: str):
    m = _QUANT_TAG_RE.match(str(text).strip())
    if m:
        return m.group("text"), m.group("tag")
    return str(text).strip(), None


def _section_strings(sec: dict) -> List[str]:
    """Every string in a section that could reach a slide."""
    out: List[str] = [str(sec.get("heading") or "")]
    body = sec.get("body") or ""
    if body:
        out.append(str(body))
    out.extend(str(b) for b in (sec.get("bullets") or []))
    table = sec.get("table")
    if isinstance(table, dict):
        out.extend(str(h) for h in (table.get("headers") or []))
        for row in table.get("rows") or []:
            out.extend(str(c) for c in row)
    for t in sec.get("tiles") or []:
        if isinstance(t, dict):
            out.extend([str(t.get("label") or ""), str(t.get("value") or "")])
    return out


def _is_empty_section(sec: dict) -> bool:
    """Drop-empty (F-60): no real content, or the docx's '(nothing logged)'
    form, or ONLY placeholder text."""
    strings = [s for s in _section_strings(sec)[1:] if s.strip()]  # skip heading
    if not strings:
        return True
    real = [
        s for s in strings
        if not _NOTHING_LOGGED_RE.search(s) and not _PLACEHOLDER_RE.search(s)
    ]
    return not real


def _check_placeholders(sec: dict) -> None:
    """Placeholder mixed into REAL content refuses the render (a dropped
    all-placeholder section never reaches here)."""
    for s in _section_strings(sec):
        if _PLACEHOLDER_RE.search(s):
            raise DeckGrammarError(
                f"placeholder text in section {sec.get('heading')!r}: {s!r} — "
                f"fill it or drop the item; the deck never ships a placeholder."
            )


# ---------------------------------------------------------------------------
# Stage 1 — the pure slide plan (grammar lives here; leak scan runs on this)
# ---------------------------------------------------------------------------

def _rows_from_section(sec: dict) -> List[str]:
    rows = [str(b).strip() for b in (sec.get("bullets") or []) if str(b).strip()]
    if not rows:
        body = str(sec.get("body") or "")
        rows = [ln.strip() for ln in body.splitlines() if ln.strip()]
    return rows


def _capped_rows(rows: List[str], cap: int):
    """(kept_rows, overflow_note_or_None) — overflow renders one honest muted
    line, never an over-cap row and never a smaller font."""
    if len(rows) <= cap:
        return rows, None
    return rows[:cap], f"+{len(rows) - cap} more in the full pack"


def build_slide_plan(
    sections: List[Dict],
    *,
    kind: str = "board_pack",
    title: str,
    subtitle: str,
    exec_header: Optional[Dict[str, str]] = None,
    asks: Optional[List[Dict[str, str]]] = None,
) -> List[dict]:
    """Sections (the SAME assembled payload make_brief rendered) → slide plan.

    Pure + deterministic + python-pptx-free. Applies the full DECK_GRAMMAR
    section→slide map, the caps, drop-empty, placeholder refusal, and the
    overflow rules. Raises DeckGrammarError on a violation.
    """
    if kind not in SUPPORTED_DECK_KINDS:
        raise ValueError(
            f"kind must be one of {sorted(SUPPORTED_DECK_KINDS)}, got {kind!r} "
            f"— the deck chokepoint is board-pack only (SPEC OUT6 scope fence)."
        )
    if not title or not str(title).strip():
        raise ValueError("title (the verdict line) is required")
    if not subtitle or not str(subtitle).strip():
        raise ValueError("subtitle (period + org) is required")
    if not isinstance(sections, list) or not sections:
        raise ValueError("sections must be a non-empty list")
    for sec in sections:
        if not isinstance(sec, dict) or not str(sec.get("heading") or "").strip():
            raise ValueError(f"each section needs a 'heading': {sec!r}")

    g = GRAMMAR
    verdict = str((exec_header or {}).get("verdict") or "").strip() or str(title).strip()
    plan: List[dict] = [{
        "slide": "title",
        "title": verdict,
        "subtitle": str(subtitle).strip(),
    }]

    ask_rows: List[str] = []
    if asks:
        if len(asks) > g["max_asks"]:
            raise DeckGrammarError(
                f"asks may hold at most {g['max_asks']} items (got {len(asks)}) "
                f"— EXEC1: more than three reader-actions is not a contract."
            )
        for a in asks:
            text = str((a or {}).get("text") or "").strip()
            if text:
                ask_rows.append(text)

    for sec in sections:
        role = _classify(sec.get("heading"))
        if role == "exec":
            continue  # the title slide IS the exec header; no agenda slide
        if _is_empty_section(sec):
            continue  # drop-empty (F-60)
        _check_placeholders(sec)
        heading = str(sec.get("heading")).strip()

        if role == "kpi":
            slide: dict = {"slide": "kpi", "heading": heading}
            tiles = sec.get("tiles")
            if tiles:
                _validate_tiles(tiles)
                slide["tiles"] = [
                    {"label": str(t["label"]).strip(), "value": str(t["value"]).strip()}
                    for t in tiles
                ]
            table = sec.get("table")
            if isinstance(table, dict) and table.get("rows"):
                # SPEC OUT3 (the OUT6 seam, now live): when a plan chart
                # renders at paint time, its PNG replaces this table beside
                # the tiles; the table STAYS IN THE PLAN as the fallback — a
                # machine with no rasterizer paints it exactly as pre-OUT3.
                slide["table"] = _plan_table(table, g["max_table_rows_per_slide"])
            # SPEC OUT3 — carry the section's chart specs into the plan
            # (validated pure, python-pptx-free; strings join the leak scan).
            # A spec that fails validation is dropped, NOT raised: the deck is
            # never stricter than the docx that carries the same payload (the
            # deck-disagrees-with-pack rule), and the refused chart's numbers
            # already live in the table. ONE chart per KPI slide — the first
            # valid spec (slide real estate; DECK_GRAMMAR §2).
            for chart in (sec.get("charts") or []):
                try:
                    validate_chart(chart)
                except (ChartDataError, ValueError):
                    continue
                slide["charts"] = [chart]
                break
            if "tiles" not in slide and "table" not in slide:
                continue  # a KPI section with neither is empty for the deck
            plan.append(slide)

        elif role in ("wins", "concerns"):
            rows, note = _capped_rows(
                _rows_from_section(sec), g["max_rows_wins_concerns"]
            )
            if not rows:
                continue
            plan.append({
                "slide": "rows",
                "heading": heading,
                "rows": [dict(zip(("text", "tag"), _split_quantify_tag(r))) for r in rows],
                "note": note,
            })

        elif role == "decisions":
            table = sec.get("table")
            if isinstance(table, dict) and table.get("rows"):
                plan.append({
                    "slide": "table",
                    "heading": heading,
                    "table": _plan_table(table, g["max_table_rows_per_slide"]),
                })
            else:
                rows, note = _capped_rows(
                    _rows_from_section(sec), g["max_table_rows_per_slide"]
                )
                if not rows:
                    continue
                plan.append({
                    "slide": "rows",
                    "heading": heading,
                    "rows": [{"text": r, "tag": None} for r in rows],
                    "note": note,
                })

        elif role == "asks":
            rows = ask_rows or _rows_from_section(sec)
            if len(rows) > g["max_asks"]:
                raise DeckGrammarError(
                    f"the Asks section carries {len(rows)} asks; the cap is "
                    f"{g['max_asks']} (EXEC1). Trim the asks — never squeeze a slide."
                )
            if not rows:
                continue
            plan.append({
                "slide": "rows",
                "heading": heading,
                "rows": [dict(zip(("text", "tag"), _split_quantify_tag(r))) for r in rows],
                "note": None,
            })
            ask_rows = []  # the section slide carried them; don't double-render

        else:  # appendix / any other content section
            plan.extend(_content_slides(sec, heading, g))

    # An asks list passed with NO asks-classified section still gets its slide.
    if ask_rows:
        plan.append({
            "slide": "rows",
            "heading": "Asks",
            "rows": [dict(zip(("text", "tag"), _split_quantify_tag(r))) for r in ask_rows],
            "note": None,
        })

    if len(plan) == 1:
        raise DeckGrammarError(
            "every section dropped as empty — there is no deck to render "
            "(the .docx pack is the deliverable; say so rather than shipping "
            "a title-only deck)."
        )
    return plan


def _plan_table(table: dict, row_cap: int) -> dict:
    rows = [[str(c) for c in r] for r in (table.get("rows") or [])]
    headers = [str(h) for h in table.get("headers")] if table.get("headers") else None
    _validate_table(rows, headers)
    note = None
    if len(rows) > row_cap:
        note = f"+{len(rows) - row_cap} more in the full pack"
        rows = rows[:row_cap]
    return {"headers": headers, "rows": rows, "note": note}


def _content_slides(sec: dict, heading: str, g: dict) -> List[dict]:
    """Appendix treatment: 1–2 content slides, ≤6 bullets each, continuation
    slide for overflow, honest note past the slide cap. A table-bearing
    appendix section renders its table as its own slide within the cap."""
    slides: List[dict] = []
    rows = _rows_from_section(sec)
    per = g["max_bullets_per_slide"]
    max_slides = g["max_content_slides_per_section"]

    table = sec.get("table")
    if isinstance(table, dict) and table.get("rows"):
        slides.append({
            "slide": "table",
            "heading": heading,
            "table": _plan_table(table, g["max_table_rows_per_slide"]),
        })

    chunks = [rows[i:i + per] for i in range(0, len(rows), per)]
    for idx, chunk in enumerate(chunks):
        if len(slides) >= max_slides:
            dropped = sum(len(c) for c in chunks[idx:])
            if slides:
                slides[-1]["note"] = f"+{dropped} more in the full pack"
            break
        slides.append({
            "slide": "rows",
            "heading": heading if idx == 0 and not slides else f"{heading} (cont.)",
            "rows": [dict(zip(("text", "tag"), _split_quantify_tag(r))) for r in chunk],
            "note": None,
        })
    return slides


def _plan_strings(plan: List[dict]) -> List[str]:
    """Every string in the plan that becomes a text run — the leak-scan input."""
    out: List[str] = []
    for s in plan:
        for key in ("title", "subtitle", "heading", "note"):
            if s.get(key):
                out.append(str(s[key]))
        for r in s.get("rows") or []:
            out.append(str(r.get("text") or ""))
            if r.get("tag"):
                out.append(str(r["tag"]))
        for t in s.get("tiles") or []:
            out.extend([str(t.get("label") or ""), str(t.get("value") or "")])
        table = s.get("table")
        if isinstance(table, dict):
            out.extend(str(h) for h in (table.get("headers") or []))
            for row in table.get("rows") or []:
                out.extend(str(c) for c in row)
            if table.get("note"):
                out.append(str(table["note"]))
        # SPEC OUT3 — chart titles/labels become PIXELS the post-save scan
        # cannot read, so they join the plan-time scan here (charts.py runs
        # its own scan too; this keeps every painted string in ONE input).
        for chart in s.get("charts") or []:
            out.extend(chart_strings(chart))
    return out


# ---------------------------------------------------------------------------
# Stage 2 — paint the plan (python-pptx; every color/font brand-resolved)
# ---------------------------------------------------------------------------

# Slide typography. Every size ≥ GRAMMAR["font_floor_pt"] — asserted below so
# the floor is a structural pin, not a convention.
_SIZES_PT = {
    "title": 40, "subtitle": 18, "heading": 28,
    "row": 18, "tag": 18, "note": 12,
    "tile_value": 30, "tile_label": 12,
    "table_header": 13, "table_cell": 12,
    "footer": 12,
}
assert min(_SIZES_PT.values()) >= GRAMMAR["font_floor_pt"], \
    "a slide size fell below the DECK_GRAMMAR font floor"

_MARGIN_IN = 0.8


def make_deck(
    output_path: str,
    sections: List[Dict],
    *,
    kind: str = "board_pack",
    title: str,
    subtitle: str,
    exec_header: Optional[Dict[str, str]] = None,
    asks: Optional[List[Dict[str, str]]] = None,
    workspace_root: Optional[str] = None,
    brand: Optional[dict] = None,
    org_id: Optional[str] = None,
) -> str:
    """Write the board-pack slide companion to `output_path` and return it.

    Args mirror `brief_writer.make_brief` where they overlap — pass the SAME
    `sections` / `exec_header` / `asks` the .docx render used (one assembly,
    two renderers; the deck can never disagree with its pack). `title` is the
    verdict line (exec_header['verdict'] wins when present); `subtitle` is
    period + org. Brand precedence: explicit `brand` dict > workspace_root /
    org_id resolution via `brand.get_brand()` > byte-stable defaults.

    Gate order: input validation + grammar (plan) → brand resolution → leak
    scan over every plan string PLUS the resolved brand footer line (the one
    painted string that does not originate in the plan) → dependency ensure →
    paint → save. Grammar and leak refusals raise BEFORE python-pptx is even
    imported — no partial file.

    Raises:
      ValueError            bad inputs / non-board_pack kind (scope fence).
      DeckGrammarError      grammar violation (placeholder, ask cap, empty deck).
      DeckLeakError         forbidden tokens in the plan — no file written.
      DeckDependencyError   python-pptx unavailable and uninstallable — the
                            caller reports it in one line and delivers the
                            .docx only. NEVER hand-build a deck around this.
    """
    plan = build_slide_plan(
        sections, kind=kind, title=title, subtitle=subtitle,
        exec_header=exec_header, asks=asks,
    )

    # Brand resolves BEFORE the scan (stdlib-only, so still ahead of the
    # python-pptx ensure): the footer line is painted on every slide but does
    # not pass through the plan, so it joins the scan input here — every
    # string that reaches the paint stage is scanned (second-eyes F-2).
    resolved = brand if brand is not None else get_brand(workspace_root, org_id)
    scan_input = _plan_strings(plan) + [str(resolved.get("footer_line") or "")]
    findings = scan_text_for_leaks("\n".join(scan_input))
    if findings:
        lines = [f"  [{f['name']}] {f['match']!r}" for f in findings[:10]]
        more = f"\n  …and {len(findings) - 10} more" if len(findings) > 10 else ""
        raise DeckLeakError(
            "Forbidden tokens in the deck plan (nothing written):\n"
            + "\n".join(lines) + more
        )

    _ensure_python_pptx()
    _paint(plan, resolved, output_path, workspace_root)

    # Detectability parity with GATE1: the .docx's gate_ran event is emitted by
    # make_brief for the same fire; the deck's audit is this stderr line plus
    # `pptx_path` on the skill's board_pack_assembled event (the event contract
    # already carries it). No second gate_ran event — the reconcile audit
    # counts one deliverable fire, not one per artifact.
    print(
        f"[deck_writer] board_pack deck rendered via make_deck — gates: "
        f"grammar, leak. path={output_path}",
        file=sys.stderr,
    )
    return str(output_path)


def _paint(plan: List[dict], brand_dict: dict, output_path: str,
           workspace_root: Optional[str]) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    p = brand_dict["palette"]
    fonts = brand_dict["fonts"]

    def rgb(hexstr: str) -> RGBColor:
        h = str(hexstr).lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    ink, heading_c, accent, muted = (
        rgb(p["ink"]), rgb(p["heading"]), rgb(p["accent"]), rgb(p["muted"]))
    tile_bg, rule_c, table_header_c, zebra = (
        rgb(p["tile_bg"]), rgb(p["rule"]), rgb(p["table_header"]), rgb(p["zebra"]))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    g = GRAMMAR
    prs = Presentation()
    prs.slide_width = Inches(g["slide_w_in"])
    prs.slide_height = Inches(g["slide_h_in"])
    blank = prs.slide_layouts[6]
    content_w = g["slide_w_in"] - 2 * _MARGIN_IN

    def run_on(para, text, *, size, color, bold=False, font=None):
        r = para.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font or fonts["body"]
        return r

    def add_footer(slide):
        tb = slide.shapes.add_textbox(
            Inches(_MARGIN_IN), Inches(g["slide_h_in"] - 0.45),
            Inches(content_w), Inches(0.3))
        para = tb.text_frame.paragraphs[0]
        run_on(para, brand_dict["footer_line"], size=_SIZES_PT["footer"], color=muted)

    def add_heading(slide, text):
        tb = slide.shapes.add_textbox(
            Inches(_MARGIN_IN), Inches(0.5), Inches(content_w), Inches(0.9))
        tb.text_frame.word_wrap = True
        para = tb.text_frame.paragraphs[0]
        run_on(para, text, size=_SIZES_PT["heading"], color=heading_c,
               bold=True, font=fonts["heading"])

    def add_note(slide, text, top_in):
        tb = slide.shapes.add_textbox(
            Inches(_MARGIN_IN), Inches(top_in), Inches(content_w), Inches(0.35))
        para = tb.text_frame.paragraphs[0]
        run_on(para, text, size=_SIZES_PT["note"], color=muted)

    def add_table(slide, spec, top_in):
        rows = spec["rows"]
        headers = spec.get("headers")
        n_rows = len(rows) + (1 if headers else 0)
        n_cols = max(len(r) for r in rows + ([headers] if headers else []))
        height = min(0.42 * n_rows, g["slide_h_in"] - top_in - 0.7)
        tbl = slide.shapes.add_table(
            n_rows, n_cols, Inches(_MARGIN_IN), Inches(top_in),
            Inches(content_w), Inches(height)).table
        # Kill the theme's default banding/first-row styling — every fill below
        # is brand-resolved (zebra), never PowerPoint's blue accent theme.
        tbl.horz_banding = False
        tbl.first_row = False
        r0 = 0
        if headers:
            for c, h in enumerate(headers):
                cell = tbl.cell(0, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = table_header_c
                para = cell.text_frame.paragraphs[0]
                run_on(para, h, size=_SIZES_PT["table_header"], color=white,
                       bold=True, font=fonts["heading"])
            r0 = 1
        for ri, row in enumerate(rows):
            for c in range(n_cols):
                val = row[c] if c < len(row) else ""
                cell = tbl.cell(ri + r0, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = zebra if ri % 2 else white
                para = cell.text_frame.paragraphs[0]
                run_on(para, val, size=_SIZES_PT["table_cell"], color=ink)
        bottom = top_in + height
        if spec.get("note"):
            add_note(slide, spec["note"], bottom + 0.1)
            bottom += 0.45
        return bottom

    for s in plan:
        slide = prs.slides.add_slide(blank)

        if s["slide"] == "title":
            logo = _resolve_logo(brand_dict.get("logo_path"), workspace_root)
            if logo:
                try:
                    slide.shapes.add_picture(
                        logo, Inches(_MARGIN_IN), Inches(0.6), height=Inches(0.5))
                except Exception:
                    pass  # a bad image never blocks the deliverable (R26)
            tb = slide.shapes.add_textbox(
                Inches(_MARGIN_IN), Inches(2.5), Inches(content_w), Inches(2.2))
            tf = tb.text_frame
            tf.word_wrap = True
            run_on(tf.paragraphs[0], s["title"], size=_SIZES_PT["title"],
                   color=heading_c, bold=True, font=fonts["heading"])
            sub = tf.add_paragraph()
            sub.space_before = Pt(14)
            run_on(sub, s["subtitle"], size=_SIZES_PT["subtitle"], color=accent)
            add_footer(slide)
            continue

        add_heading(slide, s["heading"])

        if s["slide"] == "kpi":
            top = 1.6
            tiles = s.get("tiles") or []
            if tiles:
                n = len(tiles)
                gap = 0.25
                w = (content_w - gap * (n - 1)) / n
                for i, t in enumerate(tiles):
                    shp = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(_MARGIN_IN + i * (w + gap)), Inches(top),
                        Inches(w), Inches(1.5))
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = tile_bg
                    shp.line.color.rgb = rule_c
                    shp.shadow.inherit = False
                    tf = shp.text_frame
                    tf.word_wrap = True
                    v = tf.paragraphs[0]
                    v.alignment = PP_ALIGN.CENTER
                    run_on(v, t["value"], size=_SIZES_PT["tile_value"],
                           color=heading_c, bold=True, font=fonts["heading"])
                    lab = tf.add_paragraph()
                    lab.alignment = PP_ALIGN.CENTER
                    run_on(lab, t["label"].upper(), size=_SIZES_PT["tile_label"],
                           color=accent)
                top += 1.9
            # SPEC OUT3 (the OUT6 seam): a rendered chart PNG replaces the
            # compact table beside the tiles; refusal or no rasterizer on
            # this machine falls through to the table, byte-identical to
            # pre-OUT3. Never both — one message per slide.
            chart_png = None
            for chart in s.get("charts") or []:
                try:
                    from charts import try_chart_png
                    chart_png = try_chart_png(chart, brand=brand_dict)
                except Exception:
                    chart_png = None
                if chart_png:
                    break
            if chart_png:
                try:
                    # charts.py emits 720x400 (1.8:1); fill the remaining
                    # band under the tiles, centered.
                    img_h = min(3.3, g["slide_h_in"] - top - 0.7)
                    img_w = img_h * 1.8
                    slide.shapes.add_picture(
                        chart_png,
                        Inches((g["slide_w_in"] - img_w) / 2), Inches(top),
                        height=Inches(img_h))
                except Exception:
                    chart_png = None  # a bad image never blocks (R26)
            if not chart_png and s.get("table"):
                add_table(slide, s["table"], top)

        elif s["slide"] == "rows":
            tb = slide.shapes.add_textbox(
                Inches(_MARGIN_IN), Inches(1.6), Inches(content_w),
                Inches(g["slide_h_in"] - 2.3))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, row in enumerate(s["rows"]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.space_after = Pt(10)
                run_on(para, row["text"], size=_SIZES_PT["row"], color=ink)
                if row.get("tag"):
                    run_on(para, f"   {row['tag']}", size=_SIZES_PT["tag"],
                           color=accent, bold=True)
            if s.get("note"):
                para = tf.add_paragraph()
                para.space_before = Pt(6)
                run_on(para, s["note"], size=_SIZES_PT["note"], color=muted)

        elif s["slide"] == "table":
            bottom = add_table(slide, s["table"], 1.6)
            if s.get("note"):
                add_note(slide, s["note"], bottom + 0.1)

        add_footer(slide)

    prs.save(str(output_path))


def _resolve_logo(logo_path: Optional[str], workspace_root: Optional[str]) -> Optional[str]:
    """Same posture as brief_writer._resolve_logo (R26): a configured-but-
    missing logo silently falls back to no logo — never an error in a client
    chat. Relative paths resolve against workspace_root."""
    if not logo_path:
        return None
    path = Path(logo_path)
    if not path.is_absolute() and workspace_root:
        path = Path(workspace_root) / logo_path
    try:
        return str(path) if path.is_file() else None
    except OSError:
        return None


__all__ = [
    "make_deck",
    "build_slide_plan",
    "GRAMMAR",
    "SUPPORTED_DECK_KINDS",
    "PYTHON_PPTX_PIN",
    "DeckDependencyError",
    "DeckGrammarError",
    "DeckLeakError",
]
